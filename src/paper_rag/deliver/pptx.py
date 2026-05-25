"""PowerPoint deck generator (M10 / ADR-0016).

Produces a 12-slide reading-group / academic-presentation deck for one or
more papers using ``python-pptx``. Layout is intentionally minimal — default
PowerPoint master, no custom theme, no images. The deck is meant as a
*starter* a user finishes themselves.

Slide layout (single-paper mode)
--------------------------------
1. Title
2. Outline
3. Background / Motivation
4-6. Method (3 detail slides)
7. Experiments
8. Ablations
9. Pros & Cons
10. Open Questions
11. References (paper title + arxiv link)
12. Q&A

Multi-paper mode (n>1)
----------------------
Slide structure becomes:
1. Title
2. Outline (one paper per row)
3+. Per-paper section (3 slides each: Method / Results / Critique)
   Final 3 slides: Comparison / Open Questions / References

Since python-pptx is an optional dep, we lazy-import inside ``generate`` so
the rest of paper_rag can be installed without office libs.
"""

from __future__ import annotations

import datetime as dt
import io
import logging

from ._common import (
    PaperBundle,
    collect_metadata,
    fetch_paper_bundle,
)
from .dispatch import DeliverableResult

log = logging.getLogger(__name__)


_TITLE_LAYOUT = 0          # built-in title slide
_TITLE_AND_CONTENT = 1     # built-in title + content


def _new_pres():
    """Lazy-import python-pptx; raise informative error if missing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for pptx delivery. "
            "Install with: pip install -e .[deliver]"
        ) from e
    return Presentation(), Inches, Pt


def _add_title_slide(prs, title: str, subtitle: str):
    layout = prs.slide_layouts[_TITLE_LAYOUT]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) >= 2:
        slide.placeholders[1].text = subtitle
    return slide


def _add_content_slide(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[_TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0] if bullets else ""
    for line in bullets[1:]:
        p = body.add_paragraph()
        p.text = line
    return slide


def _split_summary(summary: str) -> dict[str, list[str]]:
    """Heuristic split of a 200-word summary into method/results/limitations.

    qa_agentic was prompted to use this exact 4-section structure
    (motivation/method/results/limitations); we look for keyword anchors
    in the rendered text. If absent, fall back to chunk-by-newline.
    """
    sections = {"motivation": [], "method": [], "results": [], "limitations": []}
    current = "motivation"
    for line in (summary or "").splitlines():
        s = line.strip()
        if not s:
            continue
        low = s.lower()
        if "method" in low and len(low) < 60:
            current = "method"
            continue
        if "result" in low and len(low) < 60:
            current = "results"
            continue
        if "limit" in low and len(low) < 60:
            current = "limitations"
            continue
        sections[current].append(s)
    # Make sure each section has at least a placeholder
    for k, v in sections.items():
        if not v:
            sections[k] = [f"(no {k} described in the summary)"]
    return sections


def _build_single_paper_deck(prs, b: PaperBundle, title: str):
    subtitle = (
        f"{(b.authors[0] if b.authors else 'Unknown')}{' et al.' if len(b.authors) > 1 else ''}"
        f"{f' ({b.year})' if b.year else ''}"
    )
    _add_title_slide(prs, title or b.title or "Paper Reading", subtitle)

    _add_content_slide(prs, "Outline", [
        "Background / Motivation",
        "Method",
        "Experiments & Ablations",
        "Pros / Cons",
        "Open Questions",
    ])

    sections = _split_summary(b.summary)

    _add_content_slide(prs, "Background / Motivation", sections["motivation"][:6])

    # Three method slides — split bullets evenly
    method_bullets = sections["method"]
    chunks = _chunk_list(method_bullets, 3)
    for i, slide_bullets in enumerate(chunks, 1):
        _add_content_slide(prs, f"Method ({i}/3)", slide_bullets[:6])
    while len(chunks) < 3:
        _add_content_slide(prs, f"Method ({len(chunks)+1}/3)",
                          ["(supplementary detail — fill in from your reading)"])
        chunks.append([])

    _add_content_slide(prs, "Experiments", sections["results"][:6])

    _add_content_slide(prs, "Ablations & Notes",
                       sections["results"][6:12] or ["(extend with ablations from the paper)"])

    _add_content_slide(prs, "Pros & Cons", [
        "Pros — extracted from method/results above",
        "Cons / limitations:",
        *sections["limitations"][:4],
    ])

    _add_content_slide(prs, "Open Questions", [
        "What follow-up experiments would strengthen this work?",
        "How does it compare to the most recent baselines?",
        "Where does it fail / when not to use?",
    ])

    arxiv_str = f"arXiv:{b.arxiv_id}" if b.arxiv_id else b.paper_id
    _add_content_slide(prs, "Reference", [
        f"{b.title or 'Untitled'}",
        ", ".join(b.authors[:5]) if b.authors else "Authors unknown",
        arxiv_str,
    ])

    _add_content_slide(prs, "Q & A", ["Thanks for listening!"])


def _build_multi_paper_deck(prs, bundles: list[PaperBundle], title: str):
    _add_title_slide(prs, title or "Paper Reading Group", f"{len(bundles)} papers")

    _add_content_slide(prs, "Outline",
                       [f"{i+1}. {b.title or b.paper_id}" for i, b in enumerate(bundles[:8])])

    for i, b in enumerate(bundles, 1):
        sections = _split_summary(b.summary)
        _add_content_slide(prs, f"[{i}] Method — {b.title[:50]}",
                           sections["method"][:5])
        _add_content_slide(prs, f"[{i}] Results & Limitations",
                           (sections["results"][:3] + ["---"] + sections["limitations"][:3])[:6])
        _add_content_slide(prs, f"[{i}] Critique",
                           [
                               "Strongest claim:",
                               "Weakest claim:",
                               "How does this compare to the others?",
                           ])

    # Cross-paper summary slides
    _add_content_slide(prs, "Comparison (fill in)", [
        "Common themes:",
        "Method differences:",
        "Empirical agreement / disagreement:",
    ])
    _add_content_slide(prs, "Open Questions", [
        "Combined research gaps:",
        "Future work directions:",
    ])
    refs = []
    for b in bundles:
        arxiv_str = f"arXiv:{b.arxiv_id}" if b.arxiv_id else b.paper_id
        refs.append(f"{b.title} — {arxiv_str}")
    _add_content_slide(prs, "References", refs[:10])


def _chunk_list(items: list, n: int) -> list[list]:
    if not items:
        return []
    size = max(1, (len(items) + n - 1) // n)
    return [items[i : i + size] for i in range(0, len(items), size)][:n]


def generate(
    paper_ids: list[str],
    *,
    title: str | None = None,
    n_slides: int | None = None,  # informational; layout is fixed
    tone: str = "academic",
    **_unused,
) -> DeliverableResult:
    """Generate a .pptx deliverable."""
    title = title or "Paper Reading"

    bundles = [fetch_paper_bundle(pid) for pid in paper_ids]

    prs, _Inches, _Pt = _new_pres()

    if len(bundles) == 1:
        _build_single_paper_deck(prs, bundles[0], title)
    else:
        _build_multi_paper_deck(prs, bundles, title)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    today = dt.date.today().isoformat()
    safe_title = title.lower().replace(" ", "_").replace("/", "_")[:60]
    return DeliverableResult(
        format="pptx",
        filename=f"{safe_title}_{today}.pptx",
        content_bytes=buf.getvalue(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        metadata={
            **collect_metadata(bundles),
            "n_slides": len(prs.slides),
            "tone": tone,
        },
    )


__all__ = ["generate"]
